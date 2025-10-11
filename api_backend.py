from fastapi import FastAPI, UploadFile, File, Request, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, FileResponse
import pandas as pd
import io
import json
import os
from datetime import datetime
import re
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import seaborn as sns
import numpy as np
from sklearn.ensemble import RandomForestRegressor, GradientBoostingRegressor
from sklearn.linear_model import LinearRegression
from sklearn.metrics import r2_score, mean_squared_error, mean_absolute_error, mean_absolute_percentage_error
from sklearn.model_selection import train_test_split
import joblib
import zipfile
import pandas as pd
from docx import Document  # type: ignore
from docx.shared import Inches as DocxInches  # type: ignore
from pptx import Presentation  # type: ignore
from pptx.util import Inches, Pt  # type: ignore

app = FastAPI()

# CORS configuration: local defaults + optional env-driven origins for deployment
DEFAULT_CORS_ORIGINS = [
    "http://localhost:3000",
    "http://127.0.0.1:3000",
    "http://localhost:3001",
    "http://127.0.0.1:3001",
    "http://localhost:3010",
    "http://127.0.0.1:3010",
    "https://priceoptima-2-0.vercel.app",  # Your specific Vercel domain
    "https://priceoptima-2-0-git-main-ejikeezeani.vercel.app",  # Vercel preview domain
    "https://priceoptima-2-0-git-develop-ejikeezeani.vercel.app",  # Vercel branch domain
]

# Additional origins via env, comma-separated (e.g., https://app.example.com,https://www.example.com)
_extra_origins = os.getenv("PRICEOPTIMA_ALLOWED_ORIGINS", "")
if _extra_origins:
    DEFAULT_CORS_ORIGINS += [o.strip() for o in _extra_origins.split(",") if o.strip()]

# Allow all origins for production deployment to avoid CORS issues
_allow_all = os.getenv("PRICEOPTIMA_ALLOW_ALL_ORIGINS", "1").strip() == "1"
_cors_allow_origins = ["*"] if _allow_all else list(dict.fromkeys(DEFAULT_CORS_ORIGINS))
_cors_allow_credentials = False if _allow_all else True

app.add_middleware(
    CORSMiddleware,
    allow_origins=_cors_allow_origins,
    allow_credentials=_cors_allow_credentials,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Custom JSON encoder to handle NaN values
class NaNEncoder(json.JSONEncoder):
    def encode(self, obj):
        if isinstance(obj, float):
            if np.isnan(obj) or np.isinf(obj):
                return "0.0"
        return super().encode(obj)

def clean_nan_values(obj):
    """Recursively clean NaN and inf values from a dictionary/list"""
    if isinstance(obj, dict):
        return {key: clean_nan_values(value) for key, value in obj.items()}
    elif isinstance(obj, list):
        return [clean_nan_values(item) for item in obj]
    elif isinstance(obj, float):
        if np.isnan(obj) or np.isinf(obj):
            return 0.0
        return obj
    else:
        return obj

def generate_key_insights(data, eda_results, ml_results):
    """Generate key insights from the analysis results."""
    insights = []
    
    if data is not None:
        # Data insights
        total_records = len(data)
        insights.append(f"Dataset contains {total_records} records")
        
        if 'Revenue' in data.columns:
            total_revenue = data['Revenue'].sum()
            avg_revenue = data['Revenue'].mean()
            insights.append(f"Total revenue: ${total_revenue:,.2f}, Average: ${avg_revenue:,.2f}")
        
        if 'Product' in data.columns:
            unique_products = data['Product'].nunique()
            insights.append(f"Analysis covers {unique_products} different products")
    
    if eda_results:
        # EDA insights
        if 'correlations' in eda_results:
            corr_data = eda_results['correlations']
            if 'price_quantity' in corr_data:
                corr_val = corr_data['price_quantity']
                if abs(corr_val) > 0.5:
                    direction = "strong positive" if corr_val > 0 else "strong negative"
                    insights.append(f"Price and quantity show {direction} correlation ({corr_val:.2f})")
        
        if 'outliers' in eda_results:
            total_outliers = sum(eda_results['outliers'].values())
            if total_outliers > 0:
                insights.append(f"Detected {total_outliers} outlier values that may need attention")
    
    if ml_results:
        # ML insights
        model_type = ml_results.get('modelId', 'Unknown')
        metrics = ml_results.get('metrics', {})
        
        insights.append(f"Best performing model: {model_type}")
        
        if 'r2' in metrics:
            r2 = metrics['r2']
            if r2 > 0.8:
                insights.append(f"Excellent model performance (R² = {r2:.3f})")
            elif r2 > 0.6:
                insights.append(f"Good model performance (R² = {r2:.3f})")
            elif r2 > 0.3:
                insights.append(f"Moderate model performance (R² = {r2:.3f})")
            else:
                insights.append(f"Poor model performance (R² = {r2:.3f}) - consider feature engineering")
        
        if 'shapExplanations' in ml_results:
            shap_data = ml_results['shapExplanations']
            if 'feature_importance' in shap_data and shap_data['feature_importance']:
                top_feature = shap_data['feature_importance'][0]
                insights.append(f"Most important feature: {top_feature['feature']} (importance: {top_feature['importance']:.3f})")
        
        if 'futureForecast' in ml_results:
            forecast_data = ml_results['futureForecast']
            if 'summary' in forecast_data:
                summary = forecast_data['summary']
                trend = summary.get('trend', 'unknown')
                growth_rate = summary.get('growth_rate', 0)
                insights.append(f"Future trend: {trend} with {growth_rate:.1f}% growth rate")
    
    return insights

# Enable detailed error payloads in development when PRICEOPTIMA_DEBUG=1
DEBUG = os.getenv("PRICEOPTIMA_DEBUG", "").strip() == "1"

@app.exception_handler(Exception)
async def _unhandled_exception_handler(request: Request, exc: Exception):
    # Normalize HTTPException vs other server errors
    if isinstance(exc, HTTPException):
        content = {"error": exc.detail} if DEBUG else {"error": "Request failed"}
        return JSONResponse(status_code=exc.status_code, content=content)
    detail = {"error": "Internal Server Error"}
    if DEBUG:
        detail.update({
            "type": exc.__class__.__name__,
            "message": str(exc),
            "path": str(request.url),
        })
    return JSONResponse(status_code=500, content=detail)

# Global data storage (in production, use a database)
# Using a simple in-memory store that persists for the lifetime of the process
_data_store = {
    "uploaded_data": None,
    "processed_data": None,
    "eda_results_cache": None,
    "ml_results_cache": None,
    "eda_completed": False
}

# Backward compatibility
uploaded_data = None
processed_data = None
eda_results_cache = None
ml_results_cache = None
eda_completed = False

def get_data(key: str):
    """Get data from the store"""
    return _data_store.get(key)

def set_data(key: str, value):
    """Set data in the store"""
    _data_store[key] = value
    # Update global variables for backward compatibility
    if key == "uploaded_data":
        global uploaded_data
        uploaded_data = value
    elif key == "processed_data":
        global processed_data
        processed_data = value
    elif key == "eda_results_cache":
        global eda_results_cache
        eda_results_cache = value
    elif key == "ml_results_cache":
        global ml_results_cache
        ml_results_cache = value
    elif key == "eda_completed":
        global eda_completed
        eda_completed = value

# Constants
MAX_FILE_SIZE_BYTES = 50 * 1024 * 1024  # 50 MB
REQUIRED_COLUMNS_SYNONYMS = {
    "Date": ["date"],
    "Product": ["product", "product name", "product_name", "item", "item name"],
    "Category": ["category", "type"],
    "Price": ["price", "unit price", "unit_price"],
    "Quantity": ["quantity", "quantity sold", "qty", "qty sold"],
    "Revenue": ["revenue", "sales", "total", "total revenue"]
}
OPTIONAL_COLUMNS_SYNONYMS = {
    "Waste": ["waste", "waste amount", "waste_amount"],
    "Cost": ["cost", "unit cost", "unit_cost"],
    "Supplier": ["supplier", "vendor"]
}

# Exports directory (configurable via env var)
EXPORTS_DIR = os.getenv(
    "PRICEOPTIMA_EXPORT_DIR",
    os.path.join(os.path.expanduser("~"), "Documents", "PriceOptima", "exports"),
)
os.makedirs(EXPORTS_DIR, exist_ok=True)


def _generate_visualizations(df: pd.DataFrame) -> list[str]:
    """Generate key visualizations and return list of filenames saved in EXPORTS_DIR."""
    os.makedirs(EXPORTS_DIR, exist_ok=True)
    saved_files: list[str] = []
    try:
        # Category distribution
        if "Category" in df.columns and len(df) > 0:
            plt.figure(figsize=(8, 5))
            df["Category"].value_counts().head(15).plot(kind="bar", color="#1f77b4")
            plt.title("Category Distribution (Top 15)")
            plt.xlabel("Category")
            plt.ylabel("Count")
            plt.tight_layout()
            fn = f"category_distribution_{datetime.now().strftime('%Y%m%d_%H%M%S')}.png"
            fp = os.path.join(EXPORTS_DIR, fn)
            plt.savefig(fp)
            plt.close()
            saved_files.append(fn)

        # Sales over time (Revenue)
        if "Date" in df.columns and "Revenue" in df.columns and len(df.dropna(subset=["Date"])) > 0:
            ts_df = df.dropna(subset=["Date"]).sort_values("Date")
            plt.figure(figsize=(8, 4))
            plt.plot(ts_df["Date"], ts_df["Revenue"], color="#2ca02c")
            plt.title("Sales Over Time (Revenue)")
            plt.xlabel("Date")
            plt.ylabel("Revenue")
            plt.tight_layout()
            fn = f"sales_over_time_{datetime.now().strftime('%Y%m%d_%H%M%S')}.png"
            fp = os.path.join(EXPORTS_DIR, fn)
            plt.savefig(fp)
            plt.close()
            saved_files.append(fn)

        # Correlation heatmap
        num_cols = df.select_dtypes(include=[np.number]).columns.tolist()
        if len(num_cols) >= 2:
            plt.figure(figsize=(6, 5))
            corr = df[num_cols].corr().fillna(0)
            sns.heatmap(corr, annot=False, cmap="coolwarm", vmin=-1, vmax=1)
            plt.title("Correlation Heatmap")
            plt.tight_layout()
            fn = f"correlation_heatmap_{datetime.now().strftime('%Y%m%d_%H%M%S')}.png"
            fp = os.path.join(EXPORTS_DIR, fn)
            plt.savefig(fp)
            plt.close()
            saved_files.append(fn)
    except Exception:
        try:
            plt.close()
        except Exception:
            pass
    return saved_files


def _normalize_column_name(name: str) -> str:
    return re.sub(r"\s+", " ", str(name).strip().lower())


def standardize_columns(original_df: pd.DataFrame) -> pd.DataFrame:
    """Rename columns to standard names using synonyms. Returns a copy."""
    df = original_df.copy()
    lower_map = {i: _normalize_column_name(i) for i in df.columns}

    rename_map = {}
    # Required columns
    for standard, synonyms in REQUIRED_COLUMNS_SYNONYMS.items():
        candidates = [standard.lower(), *synonyms]
        match = next((col for col, low in lower_map.items() if low in candidates), None)
        if match is not None:
            rename_map[match] = standard

    # Optional columns
    for standard, synonyms in OPTIONAL_COLUMNS_SYNONYMS.items():
        candidates = [standard.lower(), *synonyms]
        match = next((col for col, low in lower_map.items() if low in candidates), None)
        if match is not None:
            rename_map[match] = standard

    df = df.rename(columns=rename_map)

    # Validate required columns present (after rename)
    missing = [c for c in REQUIRED_COLUMNS_SYNONYMS.keys() if c not in df.columns]
    if missing:
        raise HTTPException(status_code=400, detail=f"Missing required columns: {', '.join(missing)}")

    return df


def parse_dates(df: pd.DataFrame) -> pd.DataFrame:
    """Parse Date column supporting YYYY-MM-DD, DD/MM/YYYY, MM/DD/YYYY."""
    df = df.copy()
    if "Date" in df.columns:
        # Try multiple formats
        date_series = pd.to_datetime(df["Date"], errors="coerce", format=None)
        # If too many NaT, try dayfirst
        if date_series.isna().mean() > 0.3:
            date_series = pd.to_datetime(df["Date"], errors="coerce", dayfirst=True)
        df["Date"] = date_series
    return df


def coerce_numeric(df: pd.DataFrame, cols: list[str]) -> pd.DataFrame:
    df = df.copy()
    for col in cols:
        if col in df.columns:
            df[col] = pd.to_numeric(df[col], errors="coerce")
    return df


def preprocess_data(original_df: pd.DataFrame) -> pd.DataFrame:
    """Preprocess data: standardize columns, parse dates, coerce numerics, compute Revenue when needed, handle missing."""
    df = standardize_columns(original_df)
    df = parse_dates(df)
    df = coerce_numeric(df, ["Price", "Quantity", "Revenue", "Waste", "Cost"])

    # Compute Revenue if missing
    if "Revenue" not in df.columns or df["Revenue"].isna().all():
        if "Price" in df.columns and "Quantity" in df.columns:
            df["Revenue"] = df["Price"].fillna(0) * df["Quantity"].fillna(0)
        else:
            raise HTTPException(status_code=400, detail="Cannot compute Revenue: missing Price or Quantity")

    # Fill missing values
    for col in ["Price", "Quantity", "Revenue", "Waste", "Cost"]:
        if col in df.columns:
            df[col] = df[col].fillna(0)

    # Drop rows with missing critical categorical fields; keep rows even if Date is NaT
    # to allow non-time-based analyses to proceed.
    drop_subset = [col for col in ["Product", "Category"] if col in df.columns]
    if drop_subset:
        df = df.dropna(subset=drop_subset).reset_index(drop=True)
    return df


def compute_eda(df: pd.DataFrame) -> dict:
    """Produce summaries, distributions, correlations, outliers."""
    # Variable summaries (guard against empty datasets)
    numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
    if len(df) > 0 and numeric_cols:
        try:
            summaries = df[numeric_cols].describe().to_dict()
        except Exception:
            summaries = {}
    else:
        summaries = {}

    # Distributions
    category_distribution = df["Category"].value_counts().to_dict() if "Category" in df.columns else {}

    # Correlations
    correlations = {}
    if len(numeric_cols) >= 2:
        corr_matrix = df[numeric_cols].corr().fillna(0)
        # expose some key correlations if present
        if all(col in numeric_cols for col in ["Price", "Quantity"]):
            corr_val = corr_matrix.loc["Price", "Quantity"]
            correlations["price_quantity"] = float(corr_val) if not (np.isnan(corr_val) or np.isinf(corr_val)) else 0.0
        if all(col in numeric_cols for col in ["Price", "Revenue"]):
            corr_val = corr_matrix.loc["Price", "Revenue"]
            correlations["price_revenue"] = float(corr_val) if not (np.isnan(corr_val) or np.isinf(corr_val)) else 0.0
    
    # Outlier detection using IQR on numeric columns
    outliers = {}
    for col in numeric_cols:
        q1 = df[col].quantile(0.25)
        q3 = df[col].quantile(0.75)
        iqr = q3 - q1
        lower = q1 - 1.5 * iqr
        upper = q3 + 1.5 * iqr
        count = int(((df[col] < lower) | (df[col] > upper)).sum())
        # Handle NaN/inf values in outlier bounds
        clean_lower = float(lower) if not (np.isnan(lower) or np.isinf(lower)) else 0.0
        clean_upper = float(upper) if not (np.isnan(upper) or np.isinf(upper)) else 0.0
        outliers[col] = {"count": count, "lower": clean_lower, "upper": clean_upper}

    return {
        "overview": {
            "category_distribution": category_distribution,
            "revenue_vs_waste": {
                "revenue": df["Revenue"].tolist()[:100] if "Revenue" in df.columns else [],
                "waste": df["Waste"].tolist()[:100] if "Waste" in df.columns else []
            },
        },
        "trends": {
            # Only include rows with valid dates for time-based trend
            "sales_over_time": df.dropna(subset=["Date"]).sort_values("Date")["Revenue"].tolist()[:200] if "Date" in df.columns and "Revenue" in df.columns and len(df) > 0 else []
        },
        "correlations": correlations,
        "summaries": summaries,
        "outliers": outliers,
        "insights": [
            f"Records: {len(df):,}",
            f"Products: {df['Product'].nunique() if 'Product' in df.columns else 0}",
            f"Categories: {df['Category'].nunique() if 'Category' in df.columns else 0}",
        ],
        "recommendations": [
            "Investigate outliers for data quality issues",
            "Use category distribution to focus pricing strategy",
        ],
    }

def process_eda_data(df):
    """Process data for EDA analysis"""
    try:
        # Convert date column if it exists
        date_cols = [col for col in df.columns if 'date' in col.lower()]
        if date_cols:
            df[date_cols[0]] = pd.to_datetime(df[date_cols[0]], errors='coerce')
        
        # Get numeric columns
        numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
        
        # Calculate basic statistics
        category_col = None
        for col in df.columns:
            if 'category' in col.lower() or 'type' in col.lower():
                category_col = col
                break
        
        if category_col:
            category_distribution = df[category_col].value_counts().to_dict()
        else:
            category_distribution = {"Unknown": len(df)}
        
        # Calculate correlations
        correlations = {}
        if len(numeric_cols) >= 2:
            corr_matrix = df[numeric_cols].corr()
            correlations = {
                "price_quantity": corr_matrix.iloc[0, 1] if len(corr_matrix) > 1 else 0.5,
                "price_revenue": corr_matrix.iloc[0, 2] if len(corr_matrix) > 2 else 0.7
            }
        
        # Generate insights
        insights = [
            f"Dataset contains {len(df)} records with {len(df.columns)} features",
            f"Date range: {df[date_cols[0]].min()} to {df[date_cols[0]].max()}" if date_cols else "No date information available",
            f"Top category: {max(category_distribution, key=category_distribution.get)}" if category_distribution else "No category data"
        ]
        
        recommendations = [
            "Consider implementing dynamic pricing for high-volume categories",
            "Focus on reducing waste in perishable goods categories",
            "Analyze seasonal patterns for better inventory management"
        ]
        
        return {
            "overview": {
                "category_distribution": category_distribution,
                "revenue_vs_waste": {
                    "revenue": df[numeric_cols[0]].tolist()[:10] if numeric_cols else [1000, 1200, 900],
                    "waste": df[numeric_cols[1]].tolist()[:10] if len(numeric_cols) > 1 else [50, 60, 40]
                }
            },
            "trends": {
                "sales_over_time": df[numeric_cols[0]].tolist()[:20] if numeric_cols else [100, 120, 130, 110, 140]
            },
            "correlations": correlations,
            "insights": insights,
            "recommendations": recommendations
        }
    except Exception as e:
        return {"error": str(e)}

def train_ml_model(df, model_type="random_forest"):
    """Train ML model on the data"""
    try:
        # Prepare features
        numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
        if len(numeric_cols) < 2:
            return {"error": "Insufficient numeric data for ML training"}
        
        X = df[numeric_cols[:-1]]  # All but last column as features
        y = df[numeric_cols[-1]]   # Last column as target
        
        # Clean data of NaN and inf values
        X = X.fillna(0).replace([np.inf, -np.inf], 0)
        y = y.fillna(0).replace([np.inf, -np.inf], 0)
        
        # Select model by type
        if model_type == "random_forest":
            model = RandomForestRegressor(n_estimators=300, random_state=42, n_jobs=-1)
        elif model_type == "gradient_boosting":
            model = GradientBoostingRegressor(random_state=42)
        elif model_type == "xgboost":
            try:
                from xgboost import XGBRegressor  # type: ignore
                model = XGBRegressor(
                    n_estimators=300,
                    learning_rate=0.1,
                    max_depth=6,
                    subsample=0.9,
                    colsample_bytree=0.9,
                    random_state=42,
                    n_jobs=-1,
                )
            except Exception:
                # Fallback to Gradient Boosting if XGBoost not available
                model = GradientBoostingRegressor(random_state=42)
        elif model_type == "linear":
            model = LinearRegression()
        else:
            model = RandomForestRegressor(n_estimators=300, random_state=42, n_jobs=-1)

        # Train/test split for meaningful metrics
        X_train, X_test, y_train, y_test = train_test_split(
            X, y, test_size=0.2, random_state=42
        )
        model.fit(X_train, y_train)
        predictions = model.predict(X_test)
        
        # Calculate metrics with NaN handling
        r2 = r2_score(y_test, predictions)
        rmse = float(np.sqrt(mean_squared_error(y_test, predictions)))
        mae = float(mean_absolute_error(y_test, predictions))
        try:
            mape = float(mean_absolute_percentage_error(y_test, predictions) * 100)
        except Exception:
            mape = None
        
        # Handle NaN values for JSON serialization
        if np.isnan(r2) or np.isinf(r2):
            r2 = 0.0
        if np.isnan(rmse) or np.isinf(rmse):
            rmse = 0.0
        if np.isnan(mae) or np.isinf(mae):
            mae = 0.0
        if mape is not None and (np.isnan(mape) or np.isinf(mae)):
            mape = None
        
        # Feature importance with NaN handling
        if hasattr(model, 'feature_importances_'):
            feature_importance = []
            for col, imp in zip(X.columns, model.feature_importances_):
                # Handle NaN/inf values in feature importance
                clean_imp = float(imp) if not (np.isnan(imp) or np.isinf(imp)) else 0.0
                feature_importance.append({"feature": col, "importance": clean_imp})
        else:
            feature_importance = [
                {"feature": col, "importance": 1.0/len(X.columns)}
                for col in X.columns
            ]
        
        # Clean predictions for JSON serialization
        clean_predictions = []
        for i, (actual, pred) in enumerate(zip(y_test[:10], predictions[:10])):
            # Handle NaN/inf values in predictions
            clean_actual = float(actual) if not (np.isnan(actual) or np.isinf(actual)) else 0.0
            clean_pred = float(pred) if not (np.isnan(pred) or np.isinf(pred)) else 0.0
            clean_predictions.append({
                "actual": clean_actual, 
                "predicted": clean_pred, 
                "product": f"Sample {i+1}"
            })

        # SHAP Explainability Analysis
        shap_explanations = {}
        try:
            import shap
            # Create SHAP explainer based on model type
            if model_type == "linear":
                # For linear models, use LinearExplainer
                explainer = shap.LinearExplainer(model, X_train)
                shap_values = explainer.shap_values(X_test[:10])  # Limit to 10 samples for performance
            else:
                # For tree-based models, use TreeExplainer
                explainer = shap.TreeExplainer(model)
                shap_values = explainer.shap_values(X_test[:10])  # Limit to 10 samples for performance
            
            # Get feature names
            feature_names = list(X.columns)
            
            # Calculate mean absolute SHAP values for feature importance
            if len(shap_values.shape) == 2:  # 2D array
                mean_shap_values = np.abs(shap_values).mean(axis=0)
            else:  # 1D array
                mean_shap_values = np.abs(shap_values)
            
            # Create SHAP feature importance
            shap_importance = []
            for i, (feature, importance) in enumerate(zip(feature_names, mean_shap_values)):
                clean_importance = float(importance) if not (np.isnan(importance) or np.isinf(importance)) else 0.0
                shap_importance.append({
                    "feature": feature,
                    "importance": clean_importance
                })
            
            # Sort by importance
            shap_importance.sort(key=lambda x: x["importance"], reverse=True)
            
            # Sample SHAP values for visualization (first 5 samples, top 5 features)
            top_features = [item["feature"] for item in shap_importance[:5]]
            sample_shap_values = []
            
            for i in range(min(5, len(shap_values))):
                sample_data = {}
                for j, feature in enumerate(feature_names):
                    if feature in top_features:
                        clean_value = float(shap_values[i][j]) if not (np.isnan(shap_values[i][j]) or np.isinf(shap_values[i][j])) else 0.0
                        sample_data[feature] = clean_value
                sample_shap_values.append({
                    "sample": i + 1,
                    "shap_values": sample_data
                })
            
            shap_explanations = {
                "feature_importance": shap_importance,
                "sample_explanations": sample_shap_values,
                "summary": {
                    "total_features": len(feature_names),
                    "top_feature": shap_importance[0]["feature"] if shap_importance else None,
                    "top_importance": shap_importance[0]["importance"] if shap_importance else 0.0
                }
            }
        except Exception as e:
            # If SHAP fails, provide empty structure
            shap_explanations = {
                "feature_importance": [],
                "sample_explanations": [],
                "summary": {
                    "total_features": 0,
                    "top_feature": None,
                    "top_importance": 0.0
                },
                "error": f"SHAP analysis failed: {str(e)}"
            }

        # Future Forecasting
        future_forecast = {}
        try:
            # Generate future predictions for the next 7 periods
            n_future_periods = 7
            
            # Use the last available data point as base for forecasting
            last_features = X.iloc[-1:].values
            last_target = y.iloc[-1]
            
            # Calculate historical trends for more sophisticated forecasting
            if len(y) > 1:
                # Calculate recent trend
                recent_data = y.tail(min(10, len(y)))
                if len(recent_data) > 1:
                    trend_slope = np.polyfit(range(len(recent_data)), recent_data, 1)[0]
                    trend_percentage = (trend_slope / recent_data.mean()) * 100 if recent_data.mean() != 0 else 0
                else:
                    trend_percentage = 0
            else:
                trend_percentage = 0
            
            # Model-specific forecasting strategies
            if model_type == "linear":
                # For linear models, use coefficient-based forecasting
                if hasattr(model, 'coef_'):
                    future_predictions = []
                    for i in range(n_future_periods):
                        # Apply trend to features
                        trend_factor = 1 + (i * trend_percentage / 100) if trend_percentage != 0 else 1 + (i * 0.01)
                        adjusted_features = last_features * trend_factor
                        pred = model.predict(adjusted_features)[0]
                        clean_pred = float(pred) if not (np.isnan(pred) or np.isinf(pred)) else float(last_target)
                        future_predictions.append({
                            "period": i + 1,
                            "predicted_value": clean_pred,
                            "confidence": max(0.3, 1.0 - (i * 0.15)),
                            "method": "linear_trend"
                        })
                else:
                    future_predictions = []
            else:
                # For tree-based models, use ensemble-based forecasting
                future_predictions = []
                for i in range(n_future_periods):
                    # Multiple scenarios for better forecasting
                    scenarios = []
                    
                    # Scenario 1: Conservative (slight growth)
                    conservative_factor = 1 + (i * 0.01)
                    conservative_features = last_features * conservative_factor
                    conservative_pred = model.predict(conservative_features)[0]
                    
                    # Scenario 2: Moderate (trend-based)
                    moderate_factor = 1 + (i * max(0.01, trend_percentage / 100))
                    moderate_features = last_features * moderate_factor
                    moderate_pred = model.predict(moderate_features)[0]
                    
                    # Scenario 3: Optimistic (higher growth)
                    optimistic_factor = 1 + (i * 0.03)
                    optimistic_features = last_features * optimistic_factor
                    optimistic_pred = model.predict(optimistic_features)[0]
                    
                    # Average the scenarios
                    avg_pred = (conservative_pred + moderate_pred + optimistic_pred) / 3
                    clean_pred = float(avg_pred) if not (np.isnan(avg_pred) or np.isinf(avg_pred)) else float(last_target)
                    
                    # Calculate confidence based on scenario agreement
                    scenario_variance = np.var([conservative_pred, moderate_pred, optimistic_pred])
                    confidence = max(0.4, 1.0 - (i * 0.12) - (scenario_variance / 1000))
                    
                    future_predictions.append({
                        "period": i + 1,
                        "predicted_value": clean_pred,
                        "confidence": confidence,
                        "method": "ensemble_scenarios",
                        "scenarios": {
                            "conservative": float(conservative_pred),
                            "moderate": float(moderate_pred),
                            "optimistic": float(optimistic_pred)
                        }
                    })
            
            # Calculate summary statistics
            if future_predictions:
                values = [p["predicted_value"] for p in future_predictions]
                avg_prediction = np.mean(values)
                trend = "increasing" if values[-1] > values[0] else "decreasing" if values[-1] < values[0] else "stable"
                growth_rate = ((values[-1] - values[0]) / values[0] * 100) if values[0] != 0 else 0
            else:
                avg_prediction = 0.0
                trend = "unknown"
                growth_rate = 0.0
            
            future_forecast = {
                "predictions": future_predictions,
                "method": f"{model_type}_enhanced",
                "summary": {
                    "total_periods": n_future_periods,
                    "avg_prediction": avg_prediction,
                    "trend": trend,
                    "growth_rate": growth_rate,
                    "first_period": future_predictions[0]["predicted_value"] if future_predictions else 0.0,
                    "last_period": future_predictions[-1]["predicted_value"] if future_predictions else 0.0,
                    "confidence_avg": np.mean([p["confidence"] for p in future_predictions]) if future_predictions else 0.0
                },
                "model_info": {
                    "type": model_type,
                    "historical_trend": trend_percentage,
                    "base_value": float(last_target)
                }
            }
        except Exception as e:
            future_forecast = {
                "predictions": [],
                "method": "error",
                "summary": {
                    "total_periods": 0,
                    "avg_prediction": 0.0,
                    "trend": "unknown",
                    "growth_rate": 0.0,
                    "first_period": 0.0,
                    "last_period": 0.0,
                    "confidence_avg": 0.0
                },
                "error": f"Future forecasting failed: {str(e)}"
            }

        results = {
            "modelId": model_type,
            "metrics": {
                "r2": float(r2),
                "rmse": rmse,
                "mae": mae,
                "mape": mape,
            },
            "predictions": clean_predictions,
            "featureImportance": feature_importance,
            "shapExplanations": shap_explanations,
            "futureForecast": future_forecast
        }

        # Persist model artifact
        os.makedirs("models", exist_ok=True)
        artifact_path = os.path.join("models", f"{model_type}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.joblib")
        try:
            joblib.dump(model, artifact_path)
            results["artifactPath"] = artifact_path
        except Exception:
            results["artifactPath"] = None

        return results
    except Exception as e:
        return {"error": str(e)}

@app.post("/upload")
async def upload_data(file: UploadFile = File(...)):
    try:
        # Validate file type by extension and content type
        filename_lower = (file.filename or "").lower()
        if not filename_lower.endswith(".csv") and (file.content_type or "").lower() != "text/csv":
            raise HTTPException(status_code=400, detail="Only CSV files are accepted")

        # Read bytes and size check
        content = await file.read()
        if len(content) > MAX_FILE_SIZE_BYTES:
            raise HTTPException(status_code=400, detail="File too large. Max size is 50 MB")

        # Parse CSV
        df = pd.read_csv(io.BytesIO(content))
        
        # Preprocess (auto on upload)
        processed = preprocess_data(df)
        
        # Store data using the new data store
        set_data("uploaded_data", df)
        set_data("processed_data", processed)
        set_data("eda_results_cache", None)
        set_data("ml_results_cache", None)
        set_data("eda_completed", False)
        
        headers = list(processed.columns)
        rows = processed.to_dict(orient="records")
        
        # Calculate summary statistics
        summary = {
            "totalRecords": len(processed),
            "dateRange": f"{processed['Date'].min()} to {processed['Date'].max()}" if 'Date' in processed.columns and len(processed) else "N/A",
            "products": processed['Product'].nunique() if 'Product' in processed.columns else 0,
            "categories": processed['Category'].nunique() if 'Category' in processed.columns else 0,
            "totalRevenue": float(processed['Revenue'].sum()) if 'Revenue' in processed.columns and not (np.isnan(processed['Revenue'].sum()) or np.isinf(processed['Revenue'].sum())) else 0,
            "avgPrice": float(processed['Price'].mean()) if 'Price' in processed.columns and not (np.isnan(processed['Price'].mean()) or np.isinf(processed['Price'].mean())) else 0,
        }
        preview = rows[:5]
        
        return {
            "files": [{"name": file.filename, "size": file.size, "type": file.content_type}],
            "headers": headers,
            "rows": rows[:1000],
            "summary": summary,
            "preview": preview,
            "totalRows": len(processed)
        }
    except HTTPException as he:
        raise he
    except Exception as e:
        return JSONResponse(
            status_code=400,
            content={"error": f"Failed to process file: {str(e)}"}
        )

@app.post("/eda")
async def run_eda(request: Request):
    processed_data = get_data("processed_data")
    if processed_data is None:
        return JSONResponse(
            status_code=400,
            content={"error": "No data uploaded. Please upload data first."}
        )
    
    try:
        # Compute EDA on processed data
        eda_results = compute_eda(processed_data.copy())
        set_data("eda_results_cache", eda_results)
        set_data("eda_completed", True)
        return eda_results
    except Exception as e:
        return JSONResponse(
            status_code=500,
            content={"error": f"EDA analysis failed: {str(e)}"}
        )

@app.post("/ml")
async def run_ml(request: Request):
    processed_data = get_data("processed_data")
    eda_completed = get_data("eda_completed")
    
    if processed_data is None:
        return JSONResponse(
            status_code=400,
            content={"error": "No data uploaded. Please upload data first."}
        )
    if not eda_completed:
        return JSONResponse(
            status_code=409,
            content={"error": "EDA must be completed before ML training."}
        )
    
    try:
        body = await request.json()
        model_type = body.get("model", "random_forest")

        # Map frontend model IDs to backend model types
        model_mapping = {
            "linear_regression": "linear",
            "linear": "linear",
            "random_forest": "random_forest",
            "rf": "random_forest",
            "xgboost": "xgboost",
            "gradient_boosting": "gradient_boosting",
            "gb": "gradient_boosting",
        }

        actual_model_type = model_mapping.get(model_type, model_type)
        
        # Train ML model on the processed data
        ml_results = train_ml_model(processed_data.copy(), actual_model_type)
        
        # Debug: Check for NaN values before cleaning
        if DEBUG:
            print(f"ML results before cleaning: {ml_results}")
        
        # Clean NaN values from the response
        cleaned_results = clean_nan_values(ml_results)
        
        # Debug: Check for NaN values after cleaning
        if DEBUG:
            print(f"ML results after cleaning: {cleaned_results}")
        
        if "error" in cleaned_results:
            return JSONResponse(
                status_code=500,
                content=cleaned_results
            )
        else:
            # Cache the results for SHAP and forecasting endpoints
            set_data("ml_results_cache", cleaned_results)
            return JSONResponse(
                status_code=200,
                content=cleaned_results
            )
    except Exception as e:
        return JSONResponse(
            status_code=500,
            content={"error": f"ML training failed: {str(e)}"}
        )

@app.post("/rl")
async def run_rl(request: Request):
    body = await request.json()
    algorithm = body.get("algorithm", "dqn")
    # Return a mock RL result
    return {
        "algorithm": algorithm,
        "finalReward": 123.4,
        "convergenceEpisode": 700,
        "policy": {
            "wasteReduction": 23.5,
            "profitIncrease": 18.2,
            "customerSatisfaction": 0.87,
        },
        "trainingCurve": [
            {"episode": i * 10, "reward": 50 + i * 0.5 + (i % 5) * 2} for i in range(100)
        ],
    }

@app.post("/shap-analysis")
async def get_shap_analysis():
    """Get detailed SHAP analysis for the trained model"""
    processed_data = get_data("processed_data")
    ml_results_cache = get_data("ml_results_cache")
    
    if processed_data is None:
        raise HTTPException(status_code=400, detail="No data available for SHAP analysis")
    
    if ml_results_cache is None:
        raise HTTPException(status_code=400, detail="ML training must be completed before SHAP analysis")
    
    try:
        # Extract model type from cached results
        model_type = ml_results_cache.get("modelId", "random_forest")
        
        # Prepare features (same as in training)
        numeric_cols = processed_data.select_dtypes(include=[np.number]).columns.tolist()
        if len(numeric_cols) < 2:
            raise HTTPException(status_code=400, detail="Insufficient numeric data for SHAP analysis")
        
        X = processed_data[numeric_cols[:-1]]
        y = processed_data[numeric_cols[-1]]
        
        # Clean data
        X = X.fillna(0).replace([np.inf, -np.inf], 0)
        y = y.fillna(0).replace([np.inf, -np.inf], 0)
        
        # Train/test split
        X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2, random_state=42)
        
        # Train model (same as in ML endpoint)
        if model_type == "random_forest":
            model = RandomForestRegressor(n_estimators=300, random_state=42, n_jobs=-1)
        elif model_type == "gradient_boosting":
            model = GradientBoostingRegressor(random_state=42)
        elif model_type == "xgboost":
            try:
                from xgboost import XGBRegressor
                model = XGBRegressor(n_estimators=300, learning_rate=0.1, max_depth=6, random_state=42)
            except Exception:
                model = GradientBoostingRegressor(random_state=42)
        elif model_type == "linear":
            model = LinearRegression()
        else:
            model = RandomForestRegressor(n_estimators=300, random_state=42, n_jobs=-1)
        
        model.fit(X_train, y_train)
        
        # SHAP Analysis
        import shap
        
        try:
            if model_type == "linear":
                explainer = shap.LinearExplainer(model, X_train)
                shap_values = explainer.shap_values(X_test[:20])
            else:
                # For tree-based models, use TreeExplainer
                explainer = shap.TreeExplainer(model)
                shap_values = explainer.shap_values(X_test[:20])
        except Exception as e:
            # Fallback to KernelExplainer if TreeExplainer fails
            try:
                explainer = shap.KernelExplainer(model.predict, X_train[:50])  # Use subset for performance
                shap_values = explainer.shap_values(X_test[:10])  # Limit samples
            except Exception as e2:
                # If all SHAP methods fail, return empty results
                raise Exception(f"SHAP analysis failed: {str(e)}. Fallback also failed: {str(e2)}")
        
        # Calculate feature importance
        feature_names = list(X.columns)
        mean_shap_values = np.abs(shap_values).mean(axis=0) if len(shap_values.shape) == 2 else np.abs(shap_values)
        
        shap_importance = []
        for feature, importance in zip(feature_names, mean_shap_values):
            clean_importance = float(importance) if not (np.isnan(importance) or np.isinf(importance)) else 0.0
            shap_importance.append({
                "feature": feature,
                "importance": clean_importance
            })
        
        shap_importance.sort(key=lambda x: x["importance"], reverse=True)
        
        # Sample explanations
        sample_explanations = []
        for i in range(min(10, len(shap_values))):
            sample_data = {}
            for j, feature in enumerate(feature_names):
                clean_value = float(shap_values[i][j]) if not (np.isnan(shap_values[i][j]) or np.isinf(shap_values[i][j])) else 0.0
                sample_data[feature] = clean_value
            
            sample_explanations.append({
                "sample": i + 1,
                "actual_value": float(y_test.iloc[i]) if not (np.isnan(y_test.iloc[i]) or np.isinf(y_test.iloc[i])) else 0.0,
                "predicted_value": float(model.predict(X_test.iloc[i:i+1])[0]) if not (np.isnan(model.predict(X_test.iloc[i:i+1])[0]) or np.isinf(model.predict(X_test.iloc[i:i+1])[0])) else 0.0,
                "shap_values": sample_data
            })
        
        return {
            "feature_importance": shap_importance,
            "sample_explanations": sample_explanations,
            "summary": {
                "total_features": len(feature_names),
                "total_samples": len(sample_explanations),
                "top_feature": shap_importance[0]["feature"] if shap_importance else None,
                "top_importance": shap_importance[0]["importance"] if shap_importance else 0.0
            }
        }
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"SHAP analysis failed: {str(e)}")

@app.post("/future-forecast")
async def get_future_forecast(periods: int = 7):
    """Get future forecast predictions"""
    processed_data = get_data("processed_data")
    ml_results_cache = get_data("ml_results_cache")
    
    if processed_data is None:
        raise HTTPException(status_code=400, detail="No data available for forecasting")
    
    if ml_results_cache is None:
        raise HTTPException(status_code=400, detail="ML training must be completed before forecasting")
    
    try:
        # Extract model type from cached results
        model_type = ml_results_cache.get("modelId", "random_forest")
        
        # Prepare features
        numeric_cols = processed_data.select_dtypes(include=[np.number]).columns.tolist()
        if len(numeric_cols) < 2:
            raise HTTPException(status_code=400, detail="Insufficient numeric data for forecasting")
        
        X = processed_data[numeric_cols[:-1]]
        y = processed_data[numeric_cols[-1]]
        
        # Clean data
        X = X.fillna(0).replace([np.inf, -np.inf], 0)
        y = y.fillna(0).replace([np.inf, -np.inf], 0)
        
        # Train model
        if model_type == "random_forest":
            model = RandomForestRegressor(n_estimators=300, random_state=42, n_jobs=-1)
        elif model_type == "gradient_boosting":
            model = GradientBoostingRegressor(random_state=42)
        elif model_type == "xgboost":
            try:
                from xgboost import XGBRegressor
                model = XGBRegressor(n_estimators=300, learning_rate=0.1, max_depth=6, random_state=42)
            except Exception:
                model = GradientBoostingRegressor(random_state=42)
        elif model_type == "linear":
            model = LinearRegression()
        else:
            model = RandomForestRegressor(n_estimators=300, random_state=42, n_jobs=-1)
        
        model.fit(X, y)
        
        # Generate future predictions
        future_predictions = []
        last_features = X.iloc[-1:].values
        
        for i in range(periods):
            # Apply trend factor
            trend_factor = 1 + (i * 0.02)  # 2% growth per period
            adjusted_features = last_features * trend_factor
            
            pred = model.predict(adjusted_features)[0]
            clean_pred = float(pred) if not (np.isnan(pred) or np.isinf(pred)) else 0.0
            
            future_predictions.append({
                "period": i + 1,
                "predicted_value": clean_pred,
                "confidence": max(0.5, 1.0 - (i * 0.1))
            })
        
        # Calculate trend
        if len(future_predictions) > 1:
            trend = "increasing" if future_predictions[-1]["predicted_value"] > future_predictions[0]["predicted_value"] else "decreasing"
        else:
            trend = "stable"
        
        return {
            "predictions": future_predictions,
            "method": "trend_based",
            "summary": {
                "total_periods": periods,
                "avg_prediction": np.mean([p["predicted_value"] for p in future_predictions]),
                "trend": trend,
                "first_period": future_predictions[0]["predicted_value"] if future_predictions else 0.0,
                "last_period": future_predictions[-1]["predicted_value"] if future_predictions else 0.0
            }
        }
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Future forecasting failed: {str(e)}")

@app.post("/export")
async def export_results(request: Request):
    processed_data = get_data("processed_data")
    eda_results_cache = get_data("eda_results_cache")
    ml_results_cache = get_data("ml_results_cache")
    
    if processed_data is None:
        return JSONResponse(
            status_code=400,
            content={"error": "No data available for export. Please upload and analyze data first."}
        )
    
    try:
        body = await request.json()
        items = body.get("items", [])
        
        # Create exports directory if it doesn't exist
        os.makedirs(EXPORTS_DIR, exist_ok=True)
        
        exported_files = []
        
        for item in items:
            if item == "summary_report":
                # Generate comprehensive summary report
                report_data = {
                    "timestamp": datetime.now().isoformat(),
                    "project_info": {
                        "name": "PriceOptima Analysis Report",
                        "version": "2.0",
                        "generated_by": "PriceOptima AI System"
                    },
                    "data_summary": {
                        "total_records": len(processed_data),
                        "columns": list(processed_data.columns),
                        "date_range": f"{processed_data['Date'].min()} to {processed_data['Date'].max()}" if 'Date' in processed_data.columns and not processed_data['Date'].isna().all() else "N/A",
                        "numeric_columns": processed_data.select_dtypes(include=[np.number]).columns.tolist(),
                        "categorical_columns": processed_data.select_dtypes(include=['object']).columns.tolist(),
                        "missing_values": processed_data.isnull().sum().to_dict(),
                        "data_quality_score": (1 - processed_data.isnull().sum().sum() / (len(processed_data) * len(processed_data.columns))) * 100
                    },
                    "eda_results": eda_results_cache if eda_results_cache else {},
                    "ml_results": ml_results_cache if ml_results_cache else {},
                    "analysis_summary": {
                        "eda_completed": eda_results_cache is not None,
                        "ml_completed": ml_results_cache is not None,
                        "best_model": ml_results_cache.get("modelId", "N/A") if ml_results_cache else "N/A",
                        "model_performance": ml_results_cache.get("metrics", {}) if ml_results_cache else {},
                        "key_insights": generate_key_insights(processed_data, eda_results_cache, ml_results_cache)
                    }
                }
                # JSON summary
                json_filename = f"summary_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"
                json_path = os.path.join(EXPORTS_DIR, json_filename)
                with open(json_path, 'w', encoding='utf-8') as f:
                    json.dump(report_data, f, indent=2, ensure_ascii=False, default=str)
                exported_files.append(json_filename)

                # Enhanced DOCX report (with comprehensive analysis)
                try:
                    doc = Document()
                    doc.add_heading('PriceOptima Executive Summary', 0)
                    doc.add_paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
                    doc.add_paragraph(f"Report Version: {report_data['project_info']['version']}")
                    
                    # Executive Summary
                    doc.add_heading('Executive Summary', level=1)
                    doc.add_paragraph("This comprehensive analysis report provides insights into pricing optimization, revenue forecasting, and business intelligence using advanced machine learning techniques.")
                    
                    # Data Summary
                    doc.add_heading('Data Summary', level=1)
                    data_summary = report_data['data_summary']
                    doc.add_paragraph(f"Total Records: {data_summary['total_records']}")
                    doc.add_paragraph(f"Date Range: {data_summary['date_range']}")
                    doc.add_paragraph(f"Data Quality Score: {data_summary['data_quality_score']:.1f}%")
                    doc.add_paragraph(f"Numeric Columns: {', '.join(data_summary['numeric_columns'])}")
                    doc.add_paragraph(f"Categorical Columns: {', '.join(data_summary['categorical_columns'])}")
                    
                    # Missing Values Analysis
                    if data_summary['missing_values']:
                        doc.add_heading('Data Quality Analysis', level=2)
                        missing_table = doc.add_table(rows=1, cols=2)
                        missing_table.style = 'Table Grid'
                        hdr_cells = missing_table.rows[0].cells
                        hdr_cells[0].text = 'Column'
                        hdr_cells[1].text = 'Missing Values'
                        for col, missing_count in data_summary['missing_values'].items():
                            if missing_count > 0:
                                row_cells = missing_table.add_row().cells
                                row_cells[0].text = str(col)
                                row_cells[1].text = str(missing_count)
                    
                    # EDA Results
                    if eda_results_cache:
                        doc.add_heading('Exploratory Data Analysis', level=1)
                        doc.add_paragraph("Key findings from the exploratory data analysis:")
                        
                        if 'correlations' in eda_results_cache:
                            doc.add_heading('Correlation Analysis', level=2)
                            for corr_name, corr_value in eda_results_cache['correlations'].items():
                                doc.add_paragraph(f"{corr_name.replace('_', ' ').title()}: {corr_value:.3f}")
                        
                        if 'outliers' in eda_results_cache:
                            doc.add_heading('Outlier Detection', level=2)
                            for col, outlier_count in eda_results_cache['outliers'].items():
                                if outlier_count > 0:
                                    doc.add_paragraph(f"{col}: {outlier_count} outliers detected")
                    
                    # ML Results
                    if ml_results_cache:
                        doc.add_heading('Machine Learning Analysis', level=1)
                        model_type = ml_results_cache.get('modelId', 'Unknown')
                        metrics = ml_results_cache.get('metrics', {})
                        
                        doc.add_paragraph(f"Model Type: {model_type}")
                        doc.add_paragraph(f"R² Score: {metrics.get('r2', 'N/A')}")
                        doc.add_paragraph(f"RMSE: {metrics.get('rmse', 'N/A')}")
                        doc.add_paragraph(f"MAE: {metrics.get('mae', 'N/A')}")
                        doc.add_paragraph(f"MAPE: {metrics.get('mape', 'N/A')}%")
                        
                        # Feature Importance
                        if 'featureImportance' in ml_results_cache:
                            doc.add_heading('Feature Importance', level=2)
                            feature_table = doc.add_table(rows=1, cols=2)
                            feature_table.style = 'Table Grid'
                            hdr_cells = feature_table.rows[0].cells
                            hdr_cells[0].text = 'Feature'
                            hdr_cells[1].text = 'Importance'
                            for feature in ml_results_cache['featureImportance']:
                                row_cells = feature_table.add_row().cells
                                row_cells[0].text = feature['feature']
                                row_cells[1].text = f"{feature['importance']:.3f}"
                        
                        # SHAP Analysis
                        if 'shapExplanations' in ml_results_cache:
                            doc.add_heading('SHAP Explainability Analysis', level=2)
                            shap_data = ml_results_cache['shapExplanations']
                            if 'summary' in shap_data:
                                summary = shap_data['summary']
                                doc.add_paragraph(f"Total Features Analyzed: {summary.get('total_features', 'N/A')}")
                                doc.add_paragraph(f"Top Feature: {summary.get('top_feature', 'N/A')}")
                                doc.add_paragraph(f"Top Feature Importance: {summary.get('top_importance', 'N/A')}")
                        
                        # Future Forecasting
                        if 'futureForecast' in ml_results_cache:
                            doc.add_heading('Future Forecasting', level=2)
                            forecast_data = ml_results_cache['futureForecast']
                            if 'summary' in forecast_data:
                                summary = forecast_data['summary']
                                doc.add_paragraph(f"Forecasting Method: {forecast_data.get('method', 'N/A')}")
                                doc.add_paragraph(f"Predicted Trend: {summary.get('trend', 'N/A')}")
                                doc.add_paragraph(f"Growth Rate: {summary.get('growth_rate', 'N/A')}%")
                                doc.add_paragraph(f"Average Prediction: {summary.get('avg_prediction', 'N/A')}")
                    
                    # Key Insights
                    if 'analysis_summary' in report_data and 'key_insights' in report_data['analysis_summary']:
                        doc.add_heading('Key Insights & Recommendations', level=1)
                        insights = report_data['analysis_summary']['key_insights']
                        for insight in insights:
                            doc.add_paragraph(f"• {insight}")
                    
                    # Data Preview Table
                    doc.add_heading('Data Preview', level=1)
                    try:
                        preview = processed_data.head(10)
                        table = doc.add_table(rows=1, cols=len(preview.columns))
                        table.style = 'Table Grid'
                        hdr_cells = table.rows[0].cells
                        for i, col in enumerate(preview.columns):
                            hdr_cells[i].text = str(col)
                        for _, row in preview.iterrows():
                            row_cells = table.add_row().cells
                            for i, col in enumerate(preview.columns):
                                row_cells[i].text = str(row.get(col, ""))
                    except Exception:
                        doc.add_paragraph("Data preview not available")
                    # Add charts
                    generated = _generate_visualizations(processed_data)
                    if generated:
                        doc.add_heading('Visualizations', level=1)
                        for img_name in generated:
                            img_path = os.path.join(EXPORTS_DIR, img_name)
                            if os.path.exists(img_path):
                                doc.add_picture(img_path, width=DocxInches(5.5))
                    docx_filename = f"summary_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
                    docx_path = os.path.join(EXPORTS_DIR, docx_filename)
                    doc.save(docx_path)
                    exported_files.append(docx_filename)
                except Exception:
                    # If docx generation fails, continue with other exports
                    pass
            
            elif item == "raw_data":
                # Export processed dataset
                csv_filename = f"processed_data_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv"
                data_path = os.path.join(EXPORTS_DIR, csv_filename)
                processed_data.to_csv(data_path, index=False)
                exported_files.append(csv_filename)
                # Also export Excel (xlsx)
                try:
                    xlsx_filename = csv_filename.replace('.csv', '.xlsx')
                    xlsx_path = os.path.join(EXPORTS_DIR, xlsx_filename)
                    # Prefer openpyxl engine if available
                    processed_data.to_excel(xlsx_path, index=False, engine='openpyxl')
                    exported_files.append(xlsx_filename)
                except Exception:
                    try:
                        # Fallback to XlsxWriter
                        processed_data.to_excel(xlsx_path, index=False, engine='xlsxwriter')
                        exported_files.append(xlsx_filename)
                    except Exception:
                        pass
            
            elif item == "ml_results":
                # Generate ML results
                ml_results = train_ml_model(processed_data.copy(), "random_forest")
                ml_filename = f"ml_results_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"
                ml_path = os.path.join(EXPORTS_DIR, ml_filename)
                with open(ml_path, 'w') as f:
                    json.dump(ml_results, f, indent=2)
                exported_files.append(ml_filename)
            elif item == "eda_results":
                if eda_results_cache is None:
                    eda_now = compute_eda(processed_data.copy())
                else:
                    eda_now = eda_results_cache
                eda_filename = f"eda_results_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"
                eda_path = os.path.join(EXPORTS_DIR, eda_filename)
                with open(eda_path, 'w') as f:
                    json.dump(eda_now, f, indent=2, default=str)
                exported_files.append(eda_filename)

            elif item == "presentation":
                # Generate a richer PowerPoint presentation with visuals
                try:
                    prs = Presentation()
                    # Title slide
                    title_slide_layout = prs.slide_layouts[0]
                    slide = prs.slides.add_slide(title_slide_layout)
                    slide.shapes.title.text = "PriceOptima Analysis"
                    slide.placeholders[1].text = f"Generated {datetime.now().strftime('%Y-%m-%d %H:%M')}"

                    # Summary slide
                    summary_layout = prs.slide_layouts[1]
                    s2 = prs.slides.add_slide(summary_layout)
                    s2.shapes.title.text = "Executive Summary"
                    body = s2.placeholders[1].text_frame
                    body.clear()
                    body.text = "Key Insights"
                    recs = []
                    if eda_results_cache:
                        recs = list(eda_results_cache.get('recommendations', []))
                    elif 'eda_now' in locals() and isinstance(eda_now, dict):
                        recs = list(eda_now.get('recommendations', []))
                    if not recs:
                        recs = [
                            "Dynamic pricing improved margins.",
                            "Reduce waste in perishable categories.",
                        ]
                    for rec in recs:
                        p = body.add_paragraph()
                        p.text = str(rec)
                        p.level = 1

                    # Add a charts slide with generated images
                    imgs = _generate_visualizations(processed_data)
                    if imgs:
                        content_layout = prs.slide_layouts[5]  # Title Only
                        s3 = prs.slides.add_slide(content_layout)
                        s3.shapes.title.text = "Key Visualizations"
                        left = Inches(0.5)
                        top = Inches(1.5)
                        max_w = Inches(4.5)
                        cur_left = left
                        cur_top = top
                        for idx, img in enumerate(imgs[:4]):
                            path = os.path.join(EXPORTS_DIR, img)
                            try:
                                pic = s3.shapes.add_picture(path, cur_left, cur_top, width=max_w)
                                if idx % 2 == 1:
                                    cur_left = left
                                    cur_top = cur_top + Inches(3)
                                else:
                                    cur_left = left + max_w + Inches(0.25)
                            except Exception:
                                pass

                    pptx_filename = f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
                    pptx_path = os.path.join(EXPORTS_DIR, pptx_filename)
                    prs.save(pptx_path)
                    exported_files.append(pptx_filename)
                except Exception:
                    pass
        
        return {
            "status": "success", 
            "exported": items, 
            "files": exported_files,
            "message": f"Successfully exported {len(exported_files)} files"
        }
    except Exception as e:
        return JSONResponse(
            status_code=500,
            content={"error": f"Export failed: {str(e)}"}
        )

@app.get("/download/{filename}")
async def download_file(filename: str):
    file_path = os.path.join(EXPORTS_DIR, filename)
    if os.path.exists(file_path):
        return FileResponse(file_path, filename=filename)
    else:
        return JSONResponse(
            status_code=404,
            content={"error": "File not found"}
        )


@app.get("/download-all")
async def download_all():
    """Create a zip with all exported files and return it."""
    os.makedirs(EXPORTS_DIR, exist_ok=True)
    zip_filename = f"all_exports_{datetime.now().strftime('%Y%m%d_%H%M%S')}.zip"
    zip_path = os.path.join(EXPORTS_DIR, zip_filename)
    with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zf:
        for root, _, files in os.walk(EXPORTS_DIR):
            for f in files:
                fp = os.path.join(root, f)
                if not fp.endswith('.zip'):
                    zf.write(fp, arcname=os.path.basename(fp))
    return FileResponse(zip_path, filename=zip_filename)

@app.get("/")
async def root():
    return {"message": "PriceOptima API", "status": "running", "version": "2.0"}

@app.get("/health")
async def health_check():
    return {"status": "healthy", "message": "PriceOptima API is running"}


# Alias endpoints for compatibility with alternate frontends
@app.post("/upload-dataset")
async def upload_dataset_alias(file: UploadFile = File(...)):
    return await upload_data(file)


@app.post("/preprocess")
async def preprocess_alias():
    global processed_data
    if processed_data is None:
        raise HTTPException(status_code=400, detail="No data uploaded")
    return {"status": "ok", "rows": len(processed_data)}


@app.post("/train")
async def train_alias(request: Request):
    return await run_ml(request)


@app.get("/results")
async def results_alias():
    global processed_data, eda_results_cache
    return {
        "uploaded": processed_data is not None,
        "eda_complete": eda_completed,
        "eda": eda_results_cache,
    }


@app.get("/status")
async def status_alias():
    return {
        "uploaded": uploaded_data is not None,
        "processed": processed_data is not None,
        "eda_complete": eda_completed,
        "timestamp": datetime.now().isoformat(),
    }

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
